import logging
import os
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.util import Inches

log = logging.getLogger(__name__)


xml_duration = '''
<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
<mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">
  <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" advTm="2000" p14:dur="1000">
  </p:transition>
</mc:Choice>
<mc:Fallback>
  <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" advTm="2000">
  </p:transition>
</mc:Fallback>
</mc:AlternateContent>
'''


class Pptx:

    duration_template = "oxml_duration_template.xml"

    def __init__(self, project_name, slide_duration=20, titles=False):
        self.project_name = project_name
        self.titles = titles
        self.slide_duration = slide_duration

    def get_duration_fragment(self, duration=20):
        """
        Get the duration xml fragment to add to each slide with automatic transition
        duration: slide duration in seconds
        returns: the xml fragment to add to the slide
        """
        template_path = os.path.join(os.path.dirname(__file__), self.duration_template)
        with open(template_path) as f:
            xml_fragment = f.read()
        xml_fragment = xml_fragment.replace("{duration}", str(duration * 1000))
        print(xml_fragment)
        return parse_xml(xml_fragment)

    def generate(self, filename, images):
        # Creating presentation object
        root = Presentation()

        # Creating slide layout
        slide_layout = root.slide_layouts[0]

        """ Ref for slide types:
        0 ->  title and subtitle
        1 ->  title and content
        2 ->  section header
        3 ->  two content
        4 ->  Comparison
        5 ->  Title only
        6 ->  Blank
        7 ->  Content with caption
        8 ->  Pic with caption
        """
        if '.pptx' not in filename:
            filename = filename + '.pptx'

        path = f"generated/{self.project_name}/{filename}"
        try:
            os.makedirs(os.path.dirname(path))
        except FileExistsError:
            pass

        for title, image in images:
          slide = root.slides.add_slide(slide_layout)

          left = top = Inches(0)
          pic = slide.shapes.add_picture(image, left, top)

          if self.titles:
            slide.shapes.title.text = title

          # Adding automatic transition
          slide.element.insert(-1, self.get_duration_fragment(self.slide_duration))

        root.save(path)

        return path
