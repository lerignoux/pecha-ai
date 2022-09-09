import logging
from pptx import Presentation
from pptx.util import Inches

log = logging.getLogger(__name__)


class Ppt:
  def __init__(self, titles=False):
    self.titles = titles

  def generate(self, images, filename):
    # Creating presentation object
    root = Presentation()

    # Creating slide layout
    slide_layout = root.slide_layouts[0]

    """ Ref for slide types:
    0 ->  title and subtitle
    1 ->  title and content
    2 ->  section header
    3 ->  two content
    4 ->  Comparison
    5 ->  Title only
    6 ->  Blank
    7 ->  Content with caption
    8 ->  Pic with caption
    """

    for title, image in images:
      slide = root.slides.add_slide(slide_layout)

      if self.titles:
        slide.shapes.title.text = title

      left = top = Inches(0)
      pic = slide.shapes.add_picture(image, left, top)

      height = Inches(10)

      pic = slide.shapes.add_picture(img_path, left, top, height = height)

    root.save(filename)

    return filename
